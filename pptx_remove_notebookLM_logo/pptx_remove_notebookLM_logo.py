# /// script
# dependencies = [
#   "python-pptx",
#   "google-auth",
#   "requests",
# ]
# ///

import socket
# Force IPv4 to bypass broken IPv6 resolution on macOS
orig_getaddrinfo = socket.getaddrinfo
def forced_getaddrinfo(*args, **kwargs):
    args = list(args)
    if len(args) > 2:
        args[2] = socket.AF_INET
    else:
        while len(args) < 3:
            args.append(None)
        args[2] = socket.AF_INET
    return orig_getaddrinfo(*args, **kwargs)
socket.getaddrinfo = forced_getaddrinfo

import os
import sys
import argparse
import requests
from google.oauth2 import service_account
import google.auth.transport.requests

# Default configuration (EMUs)
RECT_WIDTH = 2200000  # ~173 pt
RECT_HEIGHT = 325000  # ~25.6 pt

def get_access_token(credentials_path):
    scopes = ['https://www.googleapis.com/auth/presentations']
    creds = service_account.Credentials.from_service_account_file(
        credentials_path, scopes=scopes)
    auth_req = google.auth.transport.requests.Request()
    creds.refresh(auth_req)
    return creds.token

def process_google_slides(presentation_id, credentials_path):
    token = get_access_token(credentials_path)
    headers = {
        'Authorization': f'Bearer {token}',
        'Content-Type': 'application/json'
    }
    
    # 1. Fetch presentation structure
    url = f'https://slides.googleapis.com/v1/presentations/{presentation_id}'
    response = requests.get(url, headers=headers)
    if response.status_code != 200:
        print(f"Error fetching presentation: {response.text}")
        sys.exit(1)
        
    presentation = response.json()
    slides = presentation.get('slides', [])
    page_size = presentation.get('pageSize', {})
    slide_width = page_size.get('width', {}).get('magnitude', 16256000)
    slide_height = page_size.get('height', {}).get('magnitude', 9144000)
    
    # Original image sizes/positions for rollback if needed
    orig_scale_x = 472.5581395348837
    orig_scale_y = 472.28125
    orig_translate_x = 0.0
    orig_translate_y = 38100.0
    
    # Position white rectangle aligned to bottom-right of image (Y = 9,105,900 EMU)
    rect_translate_x = slide_width - RECT_WIDTH
    rect_translate_y = 9105900 - RECT_HEIGHT  # 9144000 - 38100 - 650000
    
    requests_list = []
    
    for slide in slides:
        slide_id = slide.get('objectId')
        elements = slide.get('pageElements', [])
        
        # A. Rollback existing crop and transforms on images if present
        images = [el for el in elements if 'image' in el]
        for img in images:
            object_id = img.get('objectId')
            
            # Reset crop
            requests_list.append({
                'updateImageProperties': {
                    'objectId': object_id,
                    'imageProperties': {
                        'cropProperties': {
                            'leftOffset': 0.0,
                            'rightOffset': 0.0,
                            'topOffset': 0.0,
                            'bottomOffset': 0.0
                        }
                    },
                    'fields': 'cropProperties'
                }
            })
            
            # Reset transform
            requests_list.append({
                'updatePageElementTransform': {
                    'objectId': object_id,
                    'transform': {
                        'scaleX': orig_scale_x,
                        'scaleY': orig_scale_y,
                        'shearX': 0.0,
                        'shearY': 0.0,
                        'translateX': orig_translate_x,
                        'translateY': orig_translate_y,
                        'unit': 'EMU'
                    },
                    'applyMode': 'ABSOLUTE'
                }
            })
            
        # B. Check for existing cover shape to delete
        cover_id = f"cover_rect_{slide_id}"
        exists = any(el.get('objectId') == cover_id for el in elements)
        if exists:
            requests_list.append({
                'deleteObject': {
                    'objectId': cover_id
                }
            })
            
        # C. Create white cover rectangle shape
        requests_list.append({
            'createShape': {
                'objectId': cover_id,
                'shapeType': 'RECTANGLE',
                'elementProperties': {
                    'pageObjectId': slide_id,
                    'size': {
                        'width': {'magnitude': RECT_WIDTH, 'unit': 'EMU'},
                        'height': {'magnitude': RECT_HEIGHT, 'unit': 'EMU'}
                    },
                    'transform': {
                        'scaleX': 1.0,
                        'scaleY': 1.0,
                        'translateX': rect_translate_x,
                        'translateY': rect_translate_y,
                        'unit': 'EMU'
                    }
                }
            }
        })
        
        # D. Format shape with white fill, no border
        requests_list.append({
            'updateShapeProperties': {
                'objectId': cover_id,
                'shapeProperties': {
                    'shapeBackgroundFill': {
                        'solidFill': {
                            'color': {
                                'rgbColor': {
                                    'red': 1.0,
                                    'green': 1.0,
                                    'blue': 1.0
                                }
                            }
                        }
                    },
                    'outline': {
                        'propertyState': 'NOT_RENDERED'
                    }
                },
                'fields': 'shapeBackgroundFill.solidFill.color,outline'
            }
        })
        
    if not requests_list:
        print("No actions to perform.")
        return
        
    print(f"Sending batchUpdate request with {len(requests_list)} operations...")
    update_url = f'https://slides.googleapis.com/v1/presentations/{presentation_id}:batchUpdate'
    body = {'requests': requests_list}
    
    update_response = requests.post(update_url, headers=headers, json=body)
    if update_response.status_code == 200:
        print(f"Success! Processed {len(slides)} slides online.")
    else:
        print(f"BatchUpdate Error: {update_response.text}")
        sys.exit(1)

def process_local_pptx(file_path, output_path=None):
    from pptx import Presentation
    from pptx.util import Emu
    from pptx.dml.color import RGBColor

    if not os.path.exists(file_path):
        print(f"Error: File not found at {file_path}")
        sys.exit(1)

    prs = Presentation(file_path)
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    rect_width = Emu(RECT_WIDTH)
    rect_height = Emu(RECT_HEIGHT)
    rect_left = slide_width - rect_width
    
    # 38100 EMU is ~3 pt margin from the bottom
    rect_top = slide_height - rect_height - Emu(38100)
    if rect_top < 0:
        rect_top = slide_height - rect_height

    for slide in prs.slides:
        # Check if we already have a white rectangle shape at that position to avoid duplicate covers
        # (Usually covers are the latest shape, we check position as an approximation)
        shape_exists = False
        for shape in slide.shapes:
            if shape.has_chart or shape.has_table:
                continue
            if shape.shape_type == 1: # RECTANGLE
                # Check if it matches positions closely
                if abs(shape.left - rect_left) < Emu(1000) and abs(shape.top - rect_top) < Emu(1000):
                    shape_exists = True
                    break
        
        if shape_exists:
            continue

        shape = slide.shapes.add_shape(
            1, # RECTANGLE
            rect_left, rect_top, rect_width, rect_height
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(255, 255, 255)
        shape.line.fill.background()
        shape.line.width = 0

    save_path = output_path if output_path else file_path
    prs.save(save_path)
    print(f"Success! Processed {len(prs.slides)} slides locally. Saved to {save_path}")

def main():
    parser = argparse.ArgumentParser(description="Remove NotebookLM logo from slides by adding a white cover rectangle.")
    subparsers = parser.add_subparsers(dest="command", required=True, help="Command to run")
    
    # Subcommand for google-slides
    g_parser = subparsers.add_parser("google-slides", help="Process Google Slides online")
    g_parser.add_argument("presentation_id", help="The Google Slides presentation ID")
    g_parser.add_argument("--credentials", default="/Users/hoangnd/Documents/funix-auto-sheet-f464a0b5957e.json",
                        help="Path to Google Service Account Credentials JSON")
    
    # Subcommand for local-pptx
    p_parser = subparsers.add_parser("local-pptx", help="Process a local PowerPoint (.pptx) file")
    p_parser.add_argument("file_path", help="Path to local .pptx file")
    p_parser.add_argument("--output", help="Optional path to output file (defaults to overwrite input file)")
    
    args = parser.parse_args()
    
    if args.command == "google-slides":
        process_google_slides(args.presentation_id, args.credentials)
    elif args.command == "local-pptx":
        process_local_pptx(args.file_path, args.output)

if __name__ == '__main__':
    main()
